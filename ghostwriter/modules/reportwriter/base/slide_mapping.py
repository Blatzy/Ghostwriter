"""Slide mapping manager for flexible PPTX template configuration."""

import logging
from uuid import uuid4
from typing import Dict, List, Optional, Tuple

from pptx import Presentation

logger = logging.getLogger(__name__)


class SlideConfig:
    """Represents configuration for a single slide in the mapping."""

    def __init__(
        self,
        id: str,
        type: str,
        category: str,
        label: str,
        layout_index: int,
        position: int,
    ):
        self.id = id
        self.type = type
        self.category = category  # 'builtin' or 'custom'
        self.label = label
        self.layout_index = layout_index
        self.position = position

    def to_dict(self) -> dict:
        """Convert to dictionary for JSON serialization."""
        return {
            "id": self.id,
            "type": self.type,
            "category": self.category,
            "label": self.label,
            "layout_index": self.layout_index,
            "position": self.position,
        }

    @classmethod
    def from_dict(cls, data: dict) -> "SlideConfig":
        """Create from dictionary."""
        return cls(
            id=data["id"],
            type=data["type"],
            category=data["category"],
            label=data.get("label", ""),
            layout_index=data["layout_index"],
            position=data["position"],
        )


class SlideMappingManager:
    """Manages slide mapping configuration for PPTX templates."""

    # Built-in slide types with their default labels and handler method names.
    # Each built-in type has a corresponding create_* method in the exporter.
    BUILTIN_TYPES = {
        # Project slides
        "title": {"default_label": "Title Slide", "handler": "create_title_slide"},
        "agenda": {"default_label": "Agenda", "handler": "create_agenda_slide"},
        "introduction": {"default_label": "Team Introduction", "handler": "create_introduction_slide"},
        "assessment_details": {"default_label": "Assessment Details", "handler": "create_assessment_details_slide"},
        "methodology": {"default_label": "Methodology", "handler": "create_methodology_slide"},
        "timeline": {"default_label": "Assessment Timeline", "handler": "create_timeline_slide"},
        "attack_path": {"default_label": "Attack Path Overview", "handler": "create_attack_path_slide"},
        # Report slides
        "observations_overview": {"default_label": "Positive Observations Overview", "handler": "create_observations_overview_slide"},
        "observation": {"default_label": "Individual Observation Slide", "handler": "create_observation_slides"},
        "findings_overview": {"default_label": "Findings Overview", "handler": "create_findings_overview_slide"},
        "finding": {"default_label": "Individual Finding Slide", "handler": "create_finding_slides"},
        "recommendations": {"default_label": "Recommendations", "handler": "create_recommendations_slide"},
        "next_steps": {"default_label": "Next Steps", "handler": "create_next_steps_slide"},
        "final": {"default_label": "Final/Closing Slide", "handler": "create_final_slide"},
    }

    # Kept for backwards compatibility with code that references SLIDE_TYPES
    SLIDE_TYPES = {k: v["default_label"] for k, v in BUILTIN_TYPES.items()}

    # Default mapping (v2 format) for new templates
    DEFAULT_MAPPING = {
        "version": 2,
        "slides": [
            {"id": "title", "type": "title", "category": "builtin", "label": "Title Slide", "layout_index": 0, "position": 1},
            {"id": "findings_overview", "type": "findings_overview", "category": "builtin", "label": "Findings Overview", "layout_index": 1, "position": 2},
            {"id": "finding", "type": "finding", "category": "builtin", "label": "Individual Finding Slide", "layout_index": 1, "position": 3},
            {"id": "final", "type": "final", "category": "builtin", "label": "Final/Closing Slide", "layout_index": 1, "position": 4},
        ],
    }

    def __init__(
        self,
        mapping_data: Optional[dict] = None,
        presentation: Optional[Presentation] = None,
    ):
        """
        Initialize with mapping data from ReportTemplate.slide_mapping.

        Args:
            mapping_data: Dictionary from JSONField or None for defaults
            presentation: python-pptx Presentation object for validation
        """
        self.presentation = presentation

        if mapping_data is None or not isinstance(mapping_data, dict):
            logger.warning("Invalid or missing slide mapping data, using defaults")
            self.mapping_data = self.DEFAULT_MAPPING.copy()
        else:
            # Migrate v1 to v2 if needed
            if mapping_data.get("version", 1) < 2:
                self.mapping_data = self._migrate_v1_to_v2(mapping_data)
            else:
                self.mapping_data = mapping_data

        try:
            self.slides = self._parse_slides()
        except Exception as e:
            logger.exception("Failed to parse slide mapping, using defaults: %s", e)
            self.mapping_data = self.DEFAULT_MAPPING.copy()
            self.slides = self._parse_slides()

    def _migrate_v1_to_v2(self, data: dict) -> dict:
        """Convert v1 mapping format to v2."""
        new_slides = []
        for slide in data.get("slides", []):
            if not slide.get("enabled", True):
                continue  # Disabled v1 slides are dropped in v2

            stype = slide.get("type", "")
            is_builtin = stype in self.BUILTIN_TYPES

            new_slide = {
                "id": stype if is_builtin else f"custom_{uuid4().hex[:8]}",
                "type": stype if is_builtin else "custom",
                "category": "builtin" if is_builtin else "custom",
                "label": slide.get("label", "") or self.BUILTIN_TYPES.get(stype, {}).get("default_label", stype),
                "layout_index": slide.get("layout_index", 1),
                "position": slide.get("position", 1),
            }
            new_slides.append(new_slide)

        logger.info("Migrated slide mapping from v1 to v2: %d slides", len(new_slides))
        return {"version": 2, "slides": new_slides}

    def _parse_slides(self) -> List[SlideConfig]:
        """Parse slides from mapping data."""
        slides_data = self.mapping_data.get("slides", [])
        slides = []
        for s in slides_data:
            try:
                slides.append(SlideConfig.from_dict(s))
            except (KeyError, TypeError, ValueError) as e:
                logger.warning("Failed to parse slide config: %s. Skipping. Error: %s", s, e)
                continue
        return slides

    def get_slide_config(self, slide_type: str) -> Optional[SlideConfig]:
        """Get configuration for a specific built-in slide type."""
        for slide in self.slides:
            if slide.type == slide_type and slide.category == "builtin":
                return slide
        return None

    def get_layout_index(self, slide_type: str, fallback: int = 1) -> int:
        """Get layout index for a slide type with fallback."""
        config = self.get_slide_config(slide_type)
        if not config:
            return fallback

        if self.presentation:
            try:
                layout_count = len(self.presentation.slide_layouts)
                if config.layout_index >= layout_count:
                    logger.warning(
                        "Layout index %d for slide type '%s' exceeds available layouts (%d). Falling back to layout %d.",
                        config.layout_index, slide_type, layout_count, fallback,
                    )
                    return fallback
            except Exception as e:
                logger.warning("Error validating layout index: %s. Using fallback.", e)
                return fallback

        return config.layout_index

    def get_slides_by_position(self) -> List[SlideConfig]:
        """Get all slides sorted by position."""
        return sorted(self.slides, key=lambda s: s.position)

    def validate(self) -> Tuple[List[str], List[str]]:
        """
        Validate the slide mapping configuration.

        Returns:
            Tuple of (warnings, errors)
        """
        warnings = []
        errors = []

        # Check for duplicate positions
        positions = [s.position for s in self.slides]
        if len(positions) != len(set(positions)):
            warnings.append("Duplicate position values found in slide mapping")

        # Check for duplicate built-in types
        builtin_types_seen = set()
        for slide in self.slides:
            if slide.category == "builtin":
                if slide.type not in self.BUILTIN_TYPES:
                    errors.append(f"Unknown built-in type: '{slide.type}'")
                elif slide.type in builtin_types_seen:
                    errors.append(f"Duplicate built-in type: '{slide.type}'")
                else:
                    builtin_types_seen.add(slide.type)
            elif slide.category != "custom":
                errors.append(f"Invalid category '{slide.category}' for slide '{slide.label}'")

        # Validate layout indices if presentation is available
        if self.presentation:
            try:
                layout_count = len(self.presentation.slide_layouts)
                for slide in self.slides:
                    if slide.layout_index >= layout_count:
                        errors.append(
                            f"Layout index {slide.layout_index} for slide '{slide.label}' "
                            f"exceeds available layouts (0-{layout_count - 1})"
                        )
            except Exception as e:
                logger.warning("Error validating layouts: %s", e)

        return warnings, errors

    def to_dict(self) -> dict:
        """Export mapping to dictionary for JSON storage."""
        return {
            "version": 2,
            "slides": [s.to_dict() for s in self.slides],
        }

    @classmethod
    def extract_layouts_from_pptx(cls, pptx_path: str) -> List[Dict[str, any]]:
        """Extract layout information from a PPTX file."""
        try:
            prs = Presentation(pptx_path)
            layouts = []
            for idx, layout in enumerate(prs.slide_layouts):
                layouts.append({"index": idx, "name": layout.name})
            return layouts
        except Exception as e:
            logger.exception("Failed to extract layouts from %s: %s", pptx_path, e)
            return []
