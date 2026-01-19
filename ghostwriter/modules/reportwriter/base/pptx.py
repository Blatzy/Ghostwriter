
import io
import logging
import os
from datetime import date, datetime
from typing import List, Tuple

import jinja2
from django.conf import settings
from django.utils.dateformat import format as dateformat
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.parts.presentation import PresentationPart
from pptx.exc import PackageNotFoundError
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.enum.text import MSO_AUTO_SIZE

from ghostwriter.commandcenter.models import CompanyInformation
from ghostwriter.modules.reportwriter.base import ReportExportTemplateError
from ghostwriter.modules.reportwriter.base.base import ExportBase
from ghostwriter.modules.reportwriter.base.html_rich_text import LazilyRenderedTemplate
from ghostwriter.modules.reportwriter.richtext.pptx import HtmlToPptxWithEvidence
from ghostwriter.reporting.models import ReportTemplate

logger = logging.getLogger(__name__)


class ExportBasePptx(ExportBase):
    """
    Base class for exporting Pptx (PowerPoint) files

    Subclasses should override `run` to add slides to the `ppt_presentation` field, using `process_rich_text_pptx`
    to template and convert rich text fields, then return `super().run()` to save and return the presentation.
    """
    report_template: ReportTemplate
    ppt_presentation: PresentationPart
    company_config: CompanyInformation
    linting: bool

    @classmethod
    def mime_type(cls) -> str:
        return "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    @classmethod
    def extension(cls) -> str:
        return "pptx"

    def __init__(
        self,
        object,
        *,
        report_template: ReportTemplate,
        linting: bool = False,
        **kwargs
    ):
        if "jinja_debug" not in kwargs:
            kwargs["jinja_debug"] = linting
        super().__init__(object, **kwargs)
        self.linting = linting
        self.report_template = report_template

        try:
            self.ppt_presentation = Presentation(report_template.document.path)
        except PackageNotFoundError as err:
            raise ReportExportTemplateError("Template document file could not be found - try re-uploading it") from err
        except Exception:
            logger.exception(
                "Failed to load the provided template document for unknown reason: %s",
                report_template.document.path,
            )
            raise

        self.company_config = CompanyInformation.get_solo()

        # Initialize slide mapping manager
        from ghostwriter.modules.reportwriter.base.slide_mapping import SlideMappingManager
        self.slide_mapping_manager = SlideMappingManager(
            report_template.slide_mapping,
            self.ppt_presentation
        )

    def render_rich_text_pptx(self, rich_text: LazilyRenderedTemplate, slide, shape):
        """
        Renders a `LazilyRenderedTemplate`, converting the HTML from the TinyMCE rich text editor and inserting it into the passed in shape and slide.
        Converts HTML from the TinyMCE rich text editor and inserts it into the passed in slide and shape
        """
        ReportExportTemplateError.map_errors(
            lambda: HtmlToPptxWithEvidence.run(
                rich_text.render_html(),
                slide=slide,
                shape=shape,
                evidences=self.evidences_by_id,
            ),
            getattr(rich_text, "location", None)
        )

    def get_slide_context(self) -> dict:
        """
        Get Jinja2 context for static slide rendering.

        Returns:
            Dictionary with client, project, report, and other data
        """
        # Helper function to format dates
        def format_date(date_obj):
            """Format a date object to string."""
            if date_obj is None:
                return ""
            if isinstance(date_obj, str):
                return date_obj
            if hasattr(date_obj, 'strftime'):
                from django.conf import settings
                return date_obj.strftime(settings.DATE_FORMAT.replace('%', '%%').replace('%%', '%'))
            return str(date_obj)

        # Deep copy data to avoid modifying original
        import copy
        context_data = copy.deepcopy(self.data)

        # Ensure dates are formatted as strings for Jinja2
        if "project" in context_data:
            for date_field in ["start_date", "end_date"]:
                if date_field in context_data["project"]:
                    context_data["project"][date_field] = format_date(context_data["project"][date_field])

        if "report" in context_data:
            for date_field in ["complete_date", "created", "last_update"]:
                if date_field in context_data["report"]:
                    context_data["report"][date_field] = format_date(context_data["report"][date_field])

        context = {
            "client": context_data.get("client", {}),
            "project": context_data.get("project", {}),
            "report": context_data.get("report", {}),
            "team": context_data.get("team", []),
            "company": {
                "name": self.company_config.company_name,
                "email": self.company_config.company_email,
                "twitter": self.company_config.company_twitter,
            },
            "now": datetime.now(),
        }

        # Log available context for debugging
        logger.debug(
            "Jinja2 context for static slides - client: %s, project keys: %s, report keys: %s",
            context.get("client", {}).get("name", "N/A"),
            list(context.get("project", {}).keys()),
            list(context.get("report", {}).keys())
        )

        return context

    def get_slide_layout(self, slide_type: str):
        """
        Get the slide layout for a given slide type using the mapping.

        Args:
            slide_type: Type of slide (e.g., 'title', 'finding', etc.)

        Returns:
            Slide layout object, or None if slide is disabled
        """
        if not self.slide_mapping_manager.is_slide_enabled(slide_type):
            return None

        layout_index = self.slide_mapping_manager.get_layout_index(slide_type, fallback=1)
        return self.ppt_presentation.slide_layouts[layout_index]

    def get_title_shape(self, slide, shapes):
        """
        Safely get the title placeholder from a slide.

        Args:
            slide: The slide object
            shapes: The shapes collection from the slide

        Returns:
            Title shape if found, None otherwise
        """
        try:
            title = shapes.title
            if title is not None:
                return title
        except:
            pass

        # Fallback: try to find placeholder[0] (often the title)
        try:
            if 0 in shapes.placeholders:
                return shapes.placeholders[0]
        except:
            pass

        return None

    def get_placeholder(self, shapes, idx):
        """
        Safely get a placeholder by index.

        Args:
            shapes: The shapes collection from the slide
            idx: Placeholder index to retrieve

        Returns:
            Placeholder shape if found, None otherwise
        """
        try:
            if idx in shapes.placeholders:
                return shapes.placeholders[idx]
        except:
            pass
        return None

    def get_body_shape(self, slide, shapes):
        """
        Safely get the body placeholder from a slide, with fallback.

        Args:
            slide: The slide object
            shapes: The shapes collection from the slide

        Returns:
            Body shape if found, None otherwise
        """
        # Try placeholder[1] first (most common)
        placeholder = self.get_placeholder(shapes, 1)
        if placeholder:
            return placeholder

        # Fallback: find any text placeholder except title
        for shape in shapes:
            if shape.has_text_frame and shape != shapes.title:
                return shape

        return None

    def get_named_placeholder(self, shapes, name):
        """
        Get a placeholder by its name (as defined in PowerPoint master).

        Args:
            shapes: The shapes collection from the slide
            name: Name of the placeholder to find

        Returns:
            Placeholder shape if found, None otherwise
        """
        for shape in shapes.placeholders:
            if hasattr(shape, 'name') and shape.name.lower() == name.lower():
                return shape
        return None

    def process_footers(self):
        """
        Add footer elements (if there is one) to all slides based on the footer placeholder in the template
        """
        for idx, slide in enumerate(self.ppt_presentation.slides):
            date_placeholder_idx = -1
            footer_placeholder_idx = -1
            slide_number_placeholder_idx = -1
            slide_layout = slide.slide_layout

            for idx, place in enumerate(slide_layout.placeholders):
                if "Footer" in place.name:
                    footer_placeholder_idx = idx
                if "Slide Number" in place.name:
                    slide_number_placeholder_idx = idx
                if "Date" in place.name:
                    date_placeholder_idx = idx

            # Skip the title slide at index 0
            if idx > 0:
                if footer_placeholder_idx > 0:
                    footer_layout_placeholder, footer_placeholder = clone_placeholder(
                        slide, slide_layout, footer_placeholder_idx
                    )
                    footer_placeholder.text = footer_layout_placeholder.text
                if slide_number_placeholder_idx > 0:
                    _, slide_number_placeholder = clone_placeholder(
                        slide, slide_layout, slide_number_placeholder_idx
                    )
                    add_slide_number(slide_number_placeholder)
                if date_placeholder_idx > 0:
                    _, date_placeholder = clone_placeholder(
                        slide, slide_layout, date_placeholder_idx
                    )
                    date_placeholder.text = dateformat(date.today(), settings.DATE_FORMAT)

    def run(self):
        out = io.BytesIO()
        self.ppt_presentation.save(out)
        return out

    @classmethod
    def lint(cls, template_loc: str, report_template=None) -> Tuple[List[str], List[str]]:
        warnings = []
        errors = []
        try:
            if not os.path.exists(template_loc):
                logger.error("Template file path did not exist: %r", template_loc)
                errors.append("Template file does not exist – upload it again")
                return warnings, errors

            # Test 1: Check if the document is a PPTX file
            template_document = Presentation(template_loc)

            # Test 2: Check for existing slides
            slide_count = len(template_document.slides)
            logger.info("Slide count was %s", slide_count)
            if slide_count > 0:
                warnings.append(
                    "Template contains slides. If using static slide layouts, this is acceptable. "
                    "Ensure your slide master contains the layouts you want to reference."
                )

            # Test 3: Validate slide mapping if available
            if report_template and report_template.slide_mapping:
                from ghostwriter.modules.reportwriter.base.slide_mapping import SlideMappingManager

                manager = SlideMappingManager(
                    report_template.slide_mapping,
                    template_document
                )
                mapping_warnings, mapping_errors = manager.validate()
                warnings.extend(mapping_warnings)
                errors.extend(mapping_errors)

        except ReportExportTemplateError as error:
            logger.exception("Template failed linting: %s", error)
            errors.append(f"Linting failed: {error}")
        except Exception:
            logger.exception("Template failed linting")
            errors.append("Template rendering failed unexpectedly")

        logger.info("Linting finished: %d warnings, %d errors", len(warnings), len(errors))
        return warnings, errors


# Slide styles (From Master Style counting top to bottom from 0..n)
SLD_LAYOUT_TITLE = 0
SLD_LAYOUT_TITLE_AND_CONTENT = 1
SLD_LAYOUT_FINAL = 12


def add_slide_number(txtbox):
    """
    Add a slide number to the provided textbox. Ideally, the textbox should be the slide layout's slide
    number placeholder to match the template.

    Ref: https://stackoverflow.com/a/55816723
    """
    # Get a textbox's paragraph element
    par = txtbox.text_frame.paragraphs[0]._p

    # The slide number is actually a field, so we add a `fld` element to the paragraph
    # The number enclosed in the `a:t` element is the slide number and should auto-update on load/shuffle
    fld_xml = (
        '<a:fld %s id="{1F4E2DE4-8ADA-4D4E-9951-90A1D26586E7}" type="slidenum">\n'
        '  <a:rPr lang="en-US" smtClean="0"/>\n'
        "  <a:t>2</a:t>\n"
        "</a:fld>\n" % nsdecls("a")
    )
    fld = parse_xml(fld_xml)
    par.append(fld)


def clone_placeholder(slide, slide_layout, placeholder_idx):
    """
    Clone a placeholder from the slide master and return the layout and the new shape.
    """
    layout_placeholder = slide_layout.placeholders[placeholder_idx]
    slide.shapes.clone_placeholder(layout_placeholder)

    # The cloned placeholder is now the last shape in the slide
    return layout_placeholder, slide.shapes[-1]


def get_textframe(shape):
    """
    Get the shape's text frame and enable automatic resizing. The resize only
    triggers after opening the file in the PowerPoint application and making a change or saving.
    """
    text_frame = shape.text_frame
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    return text_frame


def write_bullet(text_frame, text, level):
    """Write a bullet to the provided text frame at the specified level."""
    p = text_frame.add_paragraph()
    p.text = text
    p.level = level


def write_objective_list(text_frame, objectives):
    """Write a list of objectives to the provided text frame."""
    for obj in objectives:
        status = obj["status"]
        if obj["complete"]:
            status = "Achieved"
        write_bullet(text_frame, f"{obj['objective']} – {status}", 1)


def prepare_for_pptx(value):
    """Strip HTML and clear 0x0D characters to prepare text for notes slides."""
    try:
        if value:
            return BeautifulSoup(value, "lxml").text.replace("\x0D", "")
        return "N/A"
    except Exception:
        logger.exception("Failed parsing this value for PPTX: %s", value)
        return ""


def delete_paragraph(par):
    """
    Delete the specified paragraph.

    **Parameter**

    ``par``
        Paragraph to delete from the document
    """
    p = par._p
    parent_element = p.getparent()
    if parent_element is not None:
        parent_element.remove(p)
    else:
        logger.warning("Could not delete paragraph in because it had no parent element")


def set_text_preserving_format(shape, text):
    """
    Set text in a shape while preserving the formatting from the master/layout.

    If the shape has existing text with formatting, this function preserves
    all formatting (font, size, color, etc.) and only replaces the text content.

    Args:
        shape: The shape to set text in
        text: The new text content (can be empty string)
    """
    if not shape or not shape.has_text_frame:
        return

    text_frame = shape.text_frame

    # If there are existing paragraphs with runs, preserve formatting
    if text_frame.paragraphs:
        first_para = text_frame.paragraphs[0]

        # If there are runs with formatting, use the first run's formatting
        if first_para.runs:
            first_run = first_para.runs[0]

            # Clear all runs but keep paragraph structure
            for run in list(first_para.runs):
                run.text = ""

            # Set new text in first run (preserves its formatting)
            first_run.text = text

            # Remove extra runs
            for run in list(first_para.runs[1:]):
                r = run._r
                r.getparent().remove(r)
        else:
            # No runs, just set text on paragraph (will use default formatting)
            first_para.text = text

        # Remove extra paragraphs
        for para in list(text_frame.paragraphs[1:]):
            delete_paragraph(para)
    else:
        # No paragraphs at all, set text directly
        text_frame.text = text


def add_paragraph_preserving_format(text_frame, text, level=0):
    """
    Add a paragraph while trying to preserve formatting from the layout.

    Args:
        text_frame: The text frame to add paragraph to
        text: The text content
        level: The indentation level (0-8)

    Returns:
        The created paragraph
    """
    p = text_frame.add_paragraph()
    p.text = text
    p.level = level
    return p


def copy_text_from_layout_preserving_format(dest_shape, source_ph):
    """
    Copy text and formatting from a layout placeholder to a slide shape.
    Preserves paragraphs, runs, and font properties.
    """
    if not source_ph.has_text_frame or not dest_shape.has_text_frame:
        return

    src_tf = source_ph.text_frame
    dest_tf = dest_shape.text_frame
    
    # Clear existing text
    dest_tf.clear()
    
    # If clear() leaves one empty paragraph (standard behavior), use it for the first one
    # otherwise add new paragraphs
    
    for i, src_p in enumerate(src_tf.paragraphs):
        if i == 0 and len(dest_tf.paragraphs) > 0:
            dest_p = dest_tf.paragraphs[0]
        else:
            dest_p = dest_tf.add_paragraph()
            
        # Copy paragraph properties
        dest_p.alignment = src_p.alignment
        dest_p.level = src_p.level
        
        # Copy runs
        for src_r in src_p.runs:
            dest_r = dest_p.add_run()
            dest_r.text = src_r.text
            
            # Copy font properties
            src_font = src_r.font
            dest_font = dest_r.font
            
            dest_font.name = src_font.name
            dest_font.size = src_font.size
            dest_font.bold = src_font.bold
            dest_font.italic = src_font.italic
            dest_font.underline = src_font.underline
            
            # Copy color if set
            if src_font.color and src_font.color.type:
                try:
                    if src_font.color.type == 1: # RGB
                        dest_font.color.rgb = src_font.color.rgb
                    elif src_font.color.type == 2: # THEME
                        dest_font.color.theme_color = src_font.color.theme_color
                except AttributeError:
                    pass


def render_jinja2_in_textframe(text_frame, jinja_env: jinja2.Environment, context: dict, slide, shape, evidences=None):


    """


    Render Jinja2 templates in all paragraphs of a text frame.





    Args:


        text_frame: PPTX text frame object


        jinja_env: Jinja2 environment for rendering


        context: Template context dictionary


        slide: The slide object, required for rich text processing


        shape: The shape (or cell) object containing the text_frame


        evidences: Dictionary of evidence files (optional)


    """


    if evidences is None:


        evidences = {}





    for p in text_frame.paragraphs:


        # Reconstruct the full text of the paragraph from its runs


        original_text = "".join(r.text for r in p.runs).replace('\xa0', ' ')


        


        if original_text and "{{" in original_text and "}}" in original_text:


            try:


                template = jinja_env.from_string(original_text)


                rendered_text = template.render(context)





                # Heuristic to check for leftover Jinja variables


                if '{{' in rendered_text and '}}' in rendered_text:


                    logger.warning(


                        "Jinja2 variable may not have been replaced in '%s'. Available keys: %s",


                        original_text[:100],


                        list(context.keys())


                    )





                # Simple check to see if rendered text is HTML


                is_html = bool(BeautifulSoup(rendered_text, "html.parser").find())





                # Preserve formatting of the first run


                font_props = {}


                if p.runs:


                    font = p.runs[0].font


                    font_props = {


                        'name': font.name,


                        'size': font.size,


                        'bold': font.bold,


                        'italic': font.italic,


                        'underline': font.underline,


                        'color': font.color if font.color.type else None


                    }





                # Clear all runs from the paragraph


                for r in list(p.runs):


                    p._p.remove(r._r)


                


                if is_html:


                    # Use HtmlToPptx to render the HTML


                    # Note: This will add new paragraphs, so we might need to delete the current one if it's empty


                    HtmlToPptxWithEvidence.run(rendered_text, slide, shape, evidences=evidences)


                    # If our current paragraph is now empty, we can try to remove it


                    if not p.text.strip():


                        delete_paragraph(p)


                else:


                    # Just set the text with preserved formatting


                    new_run = p.add_run()


                    new_run.text = rendered_text


                    


                    # Re-apply font properties


                    if font_props:


                        new_run.font.name = font_props.get('name')


                        new_run.font.size = font_props.get('size')


                        new_run.font.bold = font_props.get('bold')


                        new_run.font.italic = font_props.get('italic')


                        new_run.font.underline = font_props.get('underline')


                        


                        original_color = font_props.get('color')


                        if original_color:


                            try:


                                if original_color.type == 1: # RGB


                                    new_run.font.color.rgb = original_color.rgb


                                elif original_color.type == 2: # THEME


                                    new_run.font.color.theme_color = original_color.theme_color


                            except AttributeError:


                                pass





            except jinja2.exceptions.UndefinedError as e:


                logger.error(


                    "Jinja2 undefined variable in '%s': %s. Available context keys: %s",


                    original_text[:100], str(e), list(context.keys())


                )


            except Exception as e:


                logger.warning(


                    "Failed to render Jinja2 in paragraph '%s': %s",


                    original_text[:100], e


                )








def render_jinja2_in_shape(shape, jinja_env: jinja2.Environment, context: dict, slide, evidences=None):








    """








    Recursively render Jinja2 templates in a shape and its children.

















    Args:








        shape: PPTX shape object








        jinja_env: Jinja2 environment for rendering








        context: Template context dictionary








        slide: The slide object








        evidences: Dictionary of evidence files (optional)








    """








    # Handle text frames








    if shape.has_text_frame:








        # If the shape is a placeholder and empty, try to inherit text from the layout








        # This handles cases where the user put "{{ variable }}" in the Master Slide/Layout








        if shape.is_placeholder and not shape.text.strip():








            try:








                # Robustly find the layout placeholder by iterating








                # Direct access (placeholders[idx]) can fail with IndexError for high indices








                layout_ph = None








                target_idx = shape.placeholder_format.idx








                for ph in slide.slide_layout.placeholders:








                    if ph.placeholder_format.idx == target_idx:








                        layout_ph = ph








                        break








                








                if layout_ph and layout_ph.has_text_frame and layout_ph.text and "{{" in layout_ph.text:








                    # Copy the text and formatting from the layout to the slide shape








                    copy_text_from_layout_preserving_format(shape, layout_ph)

















            except (KeyError, AttributeError, Exception):








                # Ignore errors if layout placeholder mapping fails








                pass

















        render_jinja2_in_textframe(shape.text_frame, jinja_env, context, slide, shape, evidences=evidences)

















    # Handle tables








    if shape.has_table:








        for row in shape.table.rows:








            for cell in row.cells:








                render_jinja2_in_textframe(cell.text_frame, jinja_env, context, slide, cell, evidences=evidences)

















    # Handle grouped shapes








    if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP








        for child_shape in shape.shapes:








            render_jinja2_in_shape(child_shape, jinja_env, context, slide, evidences=evidences)











def create_static_slide(presentation, layout_index: int, jinja_env: jinja2.Environment, context: dict):
    """
    Create a static slide by copying a layout and rendering Jinja2 templates.

    Args:
        presentation: PPTX Presentation object
        layout_index: Index of layout to use
        jinja_env: Jinja2 environment for rendering
        context: Template context dictionary

    Returns:
        The created slide
    """
    slide_layout = presentation.slide_layouts[layout_index]
    slide = presentation.slides.add_slide(slide_layout)

    # Render Jinja2 in all shapes
    for shape in slide.shapes:
        render_jinja2_in_shape(shape, jinja_env, context, slide)

    return slide
