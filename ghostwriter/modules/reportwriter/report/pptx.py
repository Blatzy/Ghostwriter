
import io

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches
import pptx
from django.utils.safestring import mark_safe

from ghostwriter.modules.reportwriter.base.pptx import (
    ExportBasePptx,
    delete_paragraph,
    get_textframe,
    prepare_for_pptx,
    write_bullet,
    create_static_slide,
    set_text_preserving_format,
)
from ghostwriter.modules.reportwriter.project.pptx import ProjectSlidesMixin
from ghostwriter.modules.reportwriter.report.base import ExportReportBase
from ghostwriter.modules.reportwriter.richtext.pptx import HtmlToPptxWithEvidence


class ExportReportPptx(ExportBasePptx, ExportReportBase, ProjectSlidesMixin):
    def run(self) -> io.BytesIO:
        """Generate a complete PowerPoint slide deck for the current report."""

        base_context = self.map_rich_texts()
        jinja_context = self.get_slide_context()

        # Get slides sorted by position
        slides_config = self.slide_mapping_manager.get_slides_by_position()

        from ghostwriter.modules.reportwriter.base.slide_mapping import SlideMappingManager

        for slide_config in slides_config:
            if slide_config.category == "builtin":
                type_info = SlideMappingManager.BUILTIN_TYPES.get(slide_config.type)
                if type_info:
                    handler = getattr(self, type_info["handler"], None)
                    if handler:
                        handler(slide_config, base_context, jinja_context)
                    else:
                        logger.warning("No handler method '%s' for built-in type '%s'", type_info["handler"], slide_config.type)
                else:
                    logger.warning("Unknown built-in slide type: %s", slide_config.type)
            elif slide_config.category == "custom":
                create_static_slide(
                    self.ppt_presentation,
                    slide_config.layout_index,
                    self.jinja_env,
                    jinja_context,
                )

        self.process_footers()
        return super().run()

    def create_observations_overview_slide(self, config, base_context, jinja_context):
        """Create the observations overview slide."""
        slide_layout = self.ppt_presentation.slide_layouts[config.layout_index]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes

        title_shape = self.get_title_shape(slide, shapes)
        if title_shape:
            set_text_preserving_format(title_shape, "Positive Observations")

        body_shape = self.get_body_shape(slide, shapes)
        if not body_shape:
            return

        text_frame = get_textframe(body_shape)

        # If there are observations then write a table
        if len(base_context["observations"]) > 0:
            # Delete the default text placeholder
            sp = body_shape.element
            sp.getparent().remove(sp)

            # Add a table
            rows = len(base_context["observations"]) + 1
            columns = 1
            left = Inches(1.5)
            top = Inches(2)
            width = Inches(8)
            height = Inches(0.8)
            table = shapes.add_table(rows, columns, left, top, width, height).table
            # Set column width
            table.columns[0].width = Inches(10.5)
            # Write table headers
            cell = table.cell(0, 0)
            cell.text = "Observation"
            cell.fill.solid()
            cell.fill.fore_color.rgb = pptx.dml.color.RGBColor(0x2D, 0x28, 0x69)
            # Write findings rows
            row_iter = 1
            for observation in base_context["observations"]:
                table.cell(row_iter, 0).text = observation["title"]
                row_iter += 1
            # Set all cells alignment to center and vertical center
            for cell in table.iter_cells():
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        else:
            write_bullet(text_frame, "No observations", 0)

    def create_observation_slides(self, config, base_context, jinja_context):
        """
        Create individual observation slides.

        Uses Jinja2 template variables in the layout placeholders:
        - {{ title }} : Observation title
        - {{ description }} : Description of the observation

        Simply add these variables to your PowerPoint layout text and they will be replaced.
        """
        for observation in base_context["observations"]:
            slide_layout = self.ppt_presentation.slide_layouts[config.layout_index]
            observation_slide = self.ppt_presentation.slides.add_slide(slide_layout)

            # Create Jinja2 context for this specific observation
            observation_context = {
                "title": observation.get("title", ""),
                "description": mark_safe(observation.get("description", "")),
            }

            # Set the title placeholder directly (idx=0) since add_slide()
            # creates empty placeholders that don't inherit layout text
            for shape in observation_slide.placeholders:
                if shape.placeholder_format.idx == 0 and shape.has_text_frame:
                    shape.text_frame.text = observation.get("title", "")
                    break

            # Render Jinja2 variables in all shapes of the slide
            from ghostwriter.modules.reportwriter.base.pptx import render_jinja2_in_shape
            for shape in observation_slide.shapes:
                render_jinja2_in_shape(shape, self.jinja_env, observation_context, observation_slide, evidences=self.evidences_by_id)

            # Add evidence images
            for ev in observation.get("evidence", []):
                HtmlToPptxWithEvidence.make_evidence(observation_slide, ev)

    def create_findings_overview_slide(self, config, base_context, jinja_context):
        """Create the findings overview slide."""
        slide_layout = self.ppt_presentation.slide_layouts[config.layout_index]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes

        title_shape = self.get_title_shape(slide, shapes)
        if title_shape:
            set_text_preserving_format(title_shape, "Findings Overview")

        body_shape = self.get_body_shape(slide, shapes)
        if not body_shape:
            return

        text_frame = get_textframe(body_shape)

        # If there are findings then write a table of findings and severity ratings
        if len(base_context["findings"]) > 0:
            # Delete the default text placeholder
            sp = body_shape.element
            sp.getparent().remove(sp)

            # Add a table
            rows = len(base_context["findings"]) + 1
            columns = 2
            left = Inches(1.5)
            top = Inches(2)
            width = Inches(8)
            height = Inches(0.8)
            table = shapes.add_table(rows, columns, left, top, width, height).table
            # Set column width
            table.columns[0].width = Inches(8.5)
            table.columns[1].width = Inches(2.0)
            # Write table headers
            cell = table.cell(0, 0)
            cell.text = "Finding"
            cell.fill.solid()
            cell.fill.fore_color.rgb = pptx.dml.color.RGBColor(0x2D, 0x28, 0x69)
            cell = table.cell(0, 1)
            cell.text = "Severity"
            cell.fill.solid()
            cell.fill.fore_color.rgb = pptx.dml.color.RGBColor(0x2D, 0x28, 0x69)
            # Write findings rows
            row_iter = 1
            for finding in base_context["findings"]:
                table.cell(row_iter, 0).text = finding["title"]
                risk_cell = table.cell(row_iter, 1)
                # Set risk rating
                risk_cell.text = finding["severity"]
                # Set cell color fill type to solid
                risk_cell.fill.solid()
                # Color the risk cell based on corresponding severity color
                cell_color = pptx.dml.color.RGBColor(*map(lambda v: int(v, 16), finding["severity_color_hex"]))
                risk_cell.fill.fore_color.rgb = cell_color
                row_iter += 1
            # Set all cells alignment to center and vertical center
            for cell in table.iter_cells():
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        else:
            write_bullet(text_frame, "No findings", 0)

    def create_finding_slides(self, config, base_context, jinja_context):
        """
        Create individual finding slides.

        Uses Jinja2 template variables in the layout placeholders:
        - {{ title }} : Finding title
        - {{ severity }} : Severity level (Critical, High, Medium, Low)
        - {{ description }} : Description of the finding
        - {{ impact }} : Impact description
        - {{ affected_entities }} : Affected systems/entities
        - {{ mitigation }} or {{ recommendation }} : Mitigation recommendations
        - {{ replication }} or {{ replication_steps }} : Replication steps
        - {{ host_detection }} : Host detection techniques
        - {{ network_detection }} : Network detection techniques
        - {{ references }} : References

        Simply add these variables to your PowerPoint layout text and they will be replaced.
        """
        for finding in base_context["findings"]:
            slide_layout = self.ppt_presentation.slide_layouts[config.layout_index]
            finding_slide = self.ppt_presentation.slides.add_slide(slide_layout)

            # Create Jinja2 context for this specific finding
            finding_context = {
                "title": finding.get("title", ""),
                "severity": finding.get("severity", ""),
                "description": mark_safe(finding.get("description", "")),
                "impact": mark_safe(finding.get("impact", "")),
                "affected_entities": mark_safe(finding.get("affected_entities", "")),
                "mitigation": mark_safe(finding.get("recommendation", "")),
                "recommendation": mark_safe(finding.get("recommendation", "")),
                "replication": mark_safe(finding.get("replication_steps", "")),
                "replication_steps": mark_safe(finding.get("replication_steps", "")),
                "host_detection": mark_safe(finding.get("host_detection_techniques", "")),
                "network_detection": mark_safe(finding.get("network_detection_techniques", "")),
                "references": mark_safe(finding.get("references", "")),
                "cvss_score": finding.get("cvss_score", ""),
                "cvss_vector": finding.get("cvss_vector", ""),
            }

            # Set the title placeholder directly (idx=0) since add_slide()
            # creates empty placeholders that don't inherit layout text
            for shape in finding_slide.placeholders:
                if shape.placeholder_format.idx == 0 and shape.has_text_frame:
                    shape.text_frame.text = finding.get("title", "")
                    break

            # Render Jinja2 variables in all shapes of the slide
            from ghostwriter.modules.reportwriter.base.pptx import render_jinja2_in_shape
            for shape in finding_slide.shapes:
                render_jinja2_in_shape(shape, self.jinja_env, finding_context, finding_slide, evidences=self.evidences_by_id)

            # Add evidence images
            for ev in finding.get("evidence", []):
                HtmlToPptxWithEvidence.make_evidence(finding_slide, ev)

            # Add all finding data to the notes section
            entities = prepare_for_pptx(finding["affected_entities"])
            impact = prepare_for_pptx(finding["impact"])
            host_detection = prepare_for_pptx(finding["host_detection_techniques"])
            net_detection = prepare_for_pptx(finding["network_detection_techniques"])
            recommendation = prepare_for_pptx(finding["recommendation"])
            replication = prepare_for_pptx(finding["replication_steps"])
            references = prepare_for_pptx(finding["references"])
            notes_slide = finding_slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            p = text_frame.add_paragraph()
            p.text = f"""
                {finding["severity"].capitalize()}: {finding["title"]}

                AFFECTED ENTITIES
                {entities}

                IMPACT
                {impact}

                MITIGATION
                {recommendation}

                REPLICATION
                {replication}

                HOST DETECTION
                {host_detection}

                NETWORK DETECTION
                {net_detection}

                REFERENCES
                {references}
            """.replace("                ", "")

    def create_recommendations_slide(self, config, base_context, jinja_context):
        """Create recommendations slide."""
        slide_layout = self.ppt_presentation.slide_layouts[config.layout_index]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes

        title_shape = self.get_title_shape(slide, shapes)
        if title_shape:
            set_text_preserving_format(title_shape, "Recommendations")

    def create_next_steps_slide(self, config, base_context, jinja_context):
        """Create next steps slide."""
        slide_layout = self.ppt_presentation.slide_layouts[config.layout_index]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes

        title_shape = self.get_title_shape(slide, shapes)
        if title_shape:
            set_text_preserving_format(title_shape, "Next Steps")

    def create_final_slide(self, config, base_context, jinja_context):
        """Create final/closing slide."""
        slide_layout = self.ppt_presentation.slide_layouts[config.layout_index]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes

        body_shape = self.get_body_shape(slide, shapes)
        if not body_shape:
            return

        text_frame = get_textframe(body_shape)
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.line_spacing = 0.7
        p.text = self.company_config.company_name
        p = text_frame.add_paragraph()
        p.text = self.company_config.company_twitter
        p.line_spacing = 0.7
        p = text_frame.add_paragraph()
        p.text = self.company_config.company_email
        p.line_spacing = 0.7
