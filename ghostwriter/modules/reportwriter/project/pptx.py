
import io
from datetime import date

from django.conf import settings
from django.utils.dateformat import format as dateformat
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches
import pptx

from ghostwriter.modules.reportwriter.base.pptx import (
    ExportBasePptx,
    delete_paragraph,
    get_textframe,
    write_bullet,
    write_objective_list,
    create_static_slide,
    set_text_preserving_format,
)
from ghostwriter.modules.reportwriter.project.base import ExportProjectBase


class ProjectSlidesMixin:
    """
    Adds functions for generating Project-related slides - shared between the project and report exports
    """

    def create_title_slide(self, config, base_context, jinja_context):
        """Create the title slide."""
        if config.mode == "static":
            create_static_slide(self.ppt_presentation, config.layout_index, self.jinja_env, jinja_context)
            return

        slide_layout = self.ppt_presentation.slide_layouts[config.layout_index]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes

        title_shape = self.get_title_shape(slide, shapes)
        body_shape = self.get_body_shape(slide, shapes)

        if title_shape:
            set_text_preserving_format(title_shape, f'{self.data["client"]["name"]} {self.data["project"]["type"]}')

        if body_shape:
            text_frame = get_textframe(body_shape)
            # Preserve first paragraph formatting for main text
            if text_frame.paragraphs:
                text_frame.paragraphs[0].text = "Technical Outbrief"
            else:
                text_frame.text = "Technical Outbrief"
            p = text_frame.add_paragraph()
            p.text = dateformat(date.today(), settings.DATE_FORMAT)

    def create_agenda_slide(self, config, base_context, jinja_context):
        """Create the agenda slide."""
        if config.mode == "static":
            create_static_slide(self.ppt_presentation, config.layout_index, self.jinja_env, jinja_context)
            return

        slide_layout = self.ppt_presentation.slide_layouts[config.layout_index]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes

        title_shape = self.get_title_shape(slide, shapes)
        if title_shape:
            set_text_preserving_format(title_shape, "Agenda")

        body_shape = self.get_body_shape(slide, shapes)
        if not body_shape:
            return

        text_frame = get_textframe(body_shape)
        text_frame.clear()
        delete_paragraph(text_frame.paragraphs[0])

        write_bullet(text_frame, "Introduction", 0)
        write_bullet(text_frame, "Assessment Details", 0)
        write_bullet(text_frame, "Methodology", 0)
        write_bullet(text_frame, "Assessment Timeline", 0)
        write_bullet(text_frame, "Attack Path Overview", 0)
        write_bullet(text_frame, "Positive Control Observations", 0)
        write_bullet(text_frame, "Findings and Recommendations Overview", 0)
        write_bullet(text_frame, "Next Steps", 0)

    def create_introduction_slide(self, config, base_context, jinja_context):
        """Create the introduction/team slide."""
        if config.mode == "static":
            create_static_slide(self.ppt_presentation, config.layout_index, self.jinja_env, jinja_context)
            return

        slide_layout = self.ppt_presentation.slide_layouts[config.layout_index]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes

        title_shape = self.get_title_shape(slide, shapes)
        if title_shape:
            set_text_preserving_format(title_shape, "Introduction")

        body_shape = self.get_body_shape(slide, shapes)
        if not body_shape:
            return

        text_frame = get_textframe(body_shape)
        text_frame.clear()

        if self.data["team"]:
            delete_paragraph(text_frame.paragraphs[0])
            for member in self.data["team"]:
                write_bullet(text_frame, f"{member['name']} – {member['role']}", 0)
                write_bullet(text_frame, member["email"], 1)

    def create_assessment_details_slide(self, config, base_context, jinja_context):
        """Create the assessment details slide."""
        if config.mode == "static":
            create_static_slide(self.ppt_presentation, config.layout_index, self.jinja_env, jinja_context)
            return

        slide_layout = self.ppt_presentation.slide_layouts[config.layout_index]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes

        title_shape = self.get_title_shape(slide, shapes)
        if title_shape:
            set_text_preserving_format(title_shape, "Assessment Details")

        body_shape = self.get_body_shape(slide, shapes)
        if not body_shape:
            return

        text_frame = get_textframe(body_shape)
        text_frame.clear()
        delete_paragraph(text_frame.paragraphs[0])

        write_bullet(
            text_frame, f"{self.data['project']['type']} assessment of {self.data['client']['name']}", 0
        )
        write_bullet(
            text_frame,
            f"Testing performed from {self.data['project']['start_date']} to {self.data['project']['end_date']}",
            1,
        )

        self.render_rich_text_pptx(
            base_context["project"]["description_rt"],
            slide=slide,
            shape=body_shape,
        )

        # The method adds a new paragraph, so we need to get the last one to increase the indent level
        text_frame = get_textframe(body_shape)
        p = text_frame.paragraphs[-1]
        p.level = 1

        if self.data["objectives"]:
            primary_objs = []
            secondary_objs = []
            tertiary_objs = []
            for objective in self.data["objectives"]:
                if objective["priority"] == "Primary":
                    primary_objs.append(objective)
                elif objective["priority"] == "Secondary":
                    secondary_objs.append(objective)
                elif objective["priority"] == "Tertiary":
                    tertiary_objs.append(objective)

            if primary_objs:
                write_bullet(text_frame, "Primary Objectives", 0)
                write_objective_list(text_frame, primary_objs)

            if secondary_objs:
                write_bullet(text_frame, "Secondary Objectives", 0)
                write_objective_list(text_frame, secondary_objs)

            if tertiary_objs:
                write_bullet(text_frame, "Tertiary Objectives", 0)
                write_objective_list(text_frame, tertiary_objs)

    def create_methodology_slide(self, config, base_context, jinja_context):
        """Create the methodology slide."""
        if config.mode == "static":
            create_static_slide(self.ppt_presentation, config.layout_index, self.jinja_env, jinja_context)
            return

        slide_layout = self.ppt_presentation.slide_layouts[config.layout_index]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes

        title_shape = self.get_title_shape(slide, shapes)
        if title_shape:
            set_text_preserving_format(title_shape, "Methodology")

    def create_timeline_slide(self, config, base_context, jinja_context):
        """Create the assessment timeline slide."""
        if config.mode == "static":
            create_static_slide(self.ppt_presentation, config.layout_index, self.jinja_env, jinja_context)
            return

        slide_layout = self.ppt_presentation.slide_layouts[config.layout_index]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes

        title_shape = self.get_title_shape(slide, shapes)
        if title_shape:
            set_text_preserving_format(title_shape, "Assessment Timeline")

        # Delete the default text placeholder if present
        body_shape = self.get_body_shape(slide, shapes)
        if body_shape:
            sp = body_shape.element
            sp.getparent().remove(sp)

        # Add a table
        rows = 4
        columns = 2
        left = Inches(1.5)
        top = Inches(2)
        width = Inches(8)
        height = Inches(0.8)
        table = shapes.add_table(rows, columns, left, top, width, height).table
        # Set column width
        table.columns[0].width = Inches(2.0)
        table.columns[1].width = Inches(8.5)
        # Write table headers
        cell = table.cell(0, 0)
        cell.text = "Date"
        cell.fill.solid()
        cell.fill.fore_color.rgb = pptx.dml.color.RGBColor(0x2D, 0x28, 0x69)
        cell = table.cell(0, 1)
        cell.text = "Action Item"
        cell.fill.solid()
        cell.fill.fore_color.rgb = pptx.dml.color.RGBColor(0x2D, 0x28, 0x69)

        # Write date rows
        row_iter = 1
        table.cell(row_iter, 0).text = self.data["project"]["start_date"]
        table.cell(row_iter, 1).text = "Assessment execution began"
        row_iter += 1
        table.cell(row_iter, 0).text = self.data["project"]["end_date"]
        table.cell(row_iter, 1).text = "Assessment execution completed"
        row_iter += 1
        table.cell(row_iter, 0).text = self.data["project"]["end_date"]
        table.cell(row_iter, 1).text = "Draft report delivery"

        # Set all cells alignment to center and vertical center
        for cell in table.iter_cells():
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    def create_attack_path_slide(self, config, base_context, jinja_context):
        """Create the attack path overview slide."""
        if config.mode == "static":
            create_static_slide(self.ppt_presentation, config.layout_index, self.jinja_env, jinja_context)
            return

        slide_layout = self.ppt_presentation.slide_layouts[config.layout_index]
        slide = self.ppt_presentation.slides.add_slide(slide_layout)
        shapes = slide.shapes

        title_shape = self.get_title_shape(slide, shapes)
        if title_shape:
            set_text_preserving_format(title_shape, "Attack Path Overview")


class ExportProjectPptx(ExportBasePptx, ExportProjectBase, ProjectSlidesMixin):
    def run(self) -> io.BytesIO:
        base_context = self.map_rich_texts()
        jinja_context = self.get_slide_context()

        # Get enabled slides sorted by position
        slides_config = self.slide_mapping_manager.get_slides_by_position()

        # Process each slide type according to configuration
        for slide_config in slides_config:
            if not slide_config.enabled:
                continue

            slide_type = slide_config.type

            # Route to appropriate creation method
            if slide_type == "title":
                self.create_title_slide(slide_config, base_context, jinja_context)
            elif slide_type == "agenda":
                self.create_agenda_slide(slide_config, base_context, jinja_context)
            elif slide_type == "introduction":
                self.create_introduction_slide(slide_config, base_context, jinja_context)
            elif slide_type == "assessment_details":
                self.create_assessment_details_slide(slide_config, base_context, jinja_context)
            elif slide_type == "methodology":
                self.create_methodology_slide(slide_config, base_context, jinja_context)
            elif slide_type == "timeline":
                self.create_timeline_slide(slide_config, base_context, jinja_context)
            elif slide_type == "attack_path":
                self.create_attack_path_slide(slide_config, base_context, jinja_context)
            elif slide_config.mode == "static":
                # Handle custom static slides
                create_static_slide(
                    self.ppt_presentation,
                    slide_config.layout_index,
                    self.jinja_env,
                    jinja_context,
                )

        self.process_footers()
        return super().run()
